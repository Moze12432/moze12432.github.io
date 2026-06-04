"""
MozeAI Document Studio - Complete Application
Intelligent Document Workspace with AI Memory, Streaming, and Multi-Turn Context
"""

import streamlit as st
from groq import Groq
import requests
import re
import numpy as np
from bs4 import BeautifulSoup
import PyPDF2
import docx
from io import StringIO, BytesIO
import csv
import json
from datetime import datetime
import pytz
import time
import hashlib
from collections import defaultdict
import os
from difflib import unified_diff
import uuid
from typing import Dict, List, Optional, Callable, Any, Tuple
from dataclasses import dataclass, field, asdict

# ============================================================================
# DOCUMENT GENERATION LIBRARIES
# ============================================================================

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from docx import Document as WordDocument
from docx.shared import Inches as DocInches, Pt as DocPt
from docx.enum.text import WD_ALIGN_PARAGRAPH

# ============================================================================
# DATA CLASSES FOR INTELLIGENCE CORE
# ============================================================================

@dataclass
class ParsedInstruction:
    """Structured representation of user instruction"""
    intent: str  # "improve", "analyze", "transform", "generate", "create", "edit"
    target_audience: Optional[str] = None
    tone: str = "neutral"
    scope: str = "full"
    domain: str = "general"
    constraints: List[str] = field(default_factory=list)
    reasoning: str = ""
    confidence: float = 0.0
    needs_clarification: bool = False
    clarification_questions: List[str] = field(default_factory=list)
    extracted_entities: Dict[str, Any] = field(default_factory=dict)


@dataclass
class EditPlan:
    """Execution plan for document editing"""
    strategy: str
    steps: List[str] = field(default_factory=list)
    constraints: List[str] = field(default_factory=list)
    target_metrics: Dict[str, Any] = field(default_factory=dict)
    rationale: str = ""
    estimated_tokens: int = 0


@dataclass
class DocumentProfile:
    """Comprehensive document analysis result"""
    structure: Dict[str, Any] = field(default_factory=dict)
    content: Dict[str, Any] = field(default_factory=dict)
    quality: Dict[str, Any] = field(default_factory=dict)
    suggestions: List[str] = field(default_factory=list)
    strengths: List[str] = field(default_factory=list)
    metadata: Dict[str, Any] = field(default_factory=dict)


@dataclass
class EditResult:
    """Result of document edit operation"""
    edited_document: str = ""
    changes_made: Dict[str, Any] = field(default_factory=dict)
    reasoning: str = ""
    successful: bool = False
    streaming_complete: bool = True
    execution_time_ms: int = 0


@dataclass
class ConversationTurn:
    """Single conversation turn with context"""
    user_query: str = ""
    document_snapshot: Dict[str, Any] = field(default_factory=dict)
    assistant_response: str = ""
    edits_made: Dict[str, Any] = field(default_factory=dict)
    timestamp: datetime = field(default_factory=datetime.now)
    turn_id: str = field(default_factory=lambda: hashlib.md5(str(time.time()).encode()).hexdigest()[:8])


# ============================================================================
# DOCUMENT WORKSPACE CLASS
# ============================================================================

class DocumentWorkspace:
    """Manages the active document with version control and change tracking"""
    
    def __init__(self):
        self.current_document = {
            "id": str(uuid.uuid4()),
            "title": "Untitled Document",
            "content": "",
            "type": "text",
            "created_at": datetime.now().isoformat(),
            "modified_at": datetime.now().isoformat(),
            "versions": [],
            "changes": [],
            "metadata": {
                "word_count": 0,
                "char_count": 0,
                "reading_time": 0,
                "style": "general",
                "language": "english"
            }
        }
        self.version_history = []
        self.pending_changes = []
        self.suggestion_mode = False
        self.track_changes = True
        
    def update_document(self, new_content, change_description=""):
        """Update document with change tracking"""
        old_content = self.current_document["content"]
        
        if old_content == new_content:
            return False
            
        if self.track_changes:
            self.save_version(f"Before: {change_description}")
        
        changes = self._calculate_changes(old_content, new_content)
        
        self.current_document["content"] = new_content
        self.current_document["modified_at"] = datetime.now().isoformat()
        
        self._update_metadata()
        
        change_record = {
            "id": len(self.current_document["changes"]),
            "timestamp": datetime.now().isoformat(),
            "description": change_description,
            "changes": changes,
            "type": "edit"
        }
        self.current_document["changes"].append(change_record)
        
        if self.track_changes:
            self.save_version(f"After: {change_description}")
            
        return True
    
    def _calculate_changes(self, old_text, new_text):
        """Calculate specific changes between versions"""
        changes = []
        old_lines = old_text.split('\n')
        new_lines = new_text.split('\n')
        
        diff = list(unified_diff(old_lines, new_lines, lineterm=''))
        
        for line in diff:
            if line.startswith('+') and not line.startswith('+++'):
                changes.append({"type": "addition", "text": line[1:]})
            elif line.startswith('-') and not line.startswith('---'):
                changes.append({"type": "deletion", "text": line[1:]})
                
        return changes
    
    def _update_metadata(self):
        """Update document metadata"""
        content = self.current_document["content"]
        words = len(content.split())
        chars = len(content)
        
        self.current_document["metadata"]["word_count"] = words
        self.current_document["metadata"]["char_count"] = chars
        self.current_document["metadata"]["reading_time"] = max(1, words // 200)
        
    def save_version(self, description=""):
        """Save current state as version"""
        version = {
            "id": len(self.version_history),
            "timestamp": datetime.now().isoformat(),
            "content": self.current_document["content"],
            "description": description,
            "metadata": self.current_document["metadata"].copy()
        }
        self.version_history.append(version)
        
        if len(self.version_history) > 50:
            self.version_history = self.version_history[-50:]
            
        return version
    
    def restore_version(self, version_id):
        """Restore a previous version"""
        if version_id < len(self.version_history):
            version = self.version_history[version_id]
            self.update_document(version["content"], f"Restored version {version_id}")
            return True
        return False
    
    def analyze_document(self):
        """Perform comprehensive document analysis"""
        content = self.current_document["content"]
        
        analysis = {
            "structure": self._analyze_structure(),
            "readability": self._analyze_readability(),
            "grammar_issues": self._check_grammar(),
            "style_analysis": self._analyze_style(),
            "suggestions": self._generate_suggestions()
        }
        
        return analysis
    
    def _analyze_structure(self):
        """Analyze document structure"""
        content = self.current_document["content"]
        lines = content.split('\n')
        
        headings = []
        paragraphs = 0
        lists = 0
        
        for line in lines:
            if line.strip().startswith('#'):
                headings.append(line.strip())
            elif len(line.strip()) > 20:
                paragraphs += 1
            elif line.strip().startswith(('-', '*', '•')):
                lists += 1
                
        return {
            "headings": headings,
            "paragraph_count": paragraphs,
            "list_items": lists,
            "total_lines": len(lines)
        }
    
    def _analyze_readability(self):
        """Calculate readability scores"""
        content = self.current_document["content"]
        sentences = re.split(r'[.!?]+', content)
        words = content.split()
        
        if len(sentences) == 0 or len(words) == 0:
            return {"score": 0, "level": "Unknown"}
        
        avg_words_per_sentence = len(words) / len(sentences)
        
        if avg_words_per_sentence < 10:
            score = 90
            level = "Very Easy"
        elif avg_words_per_sentence < 15:
            score = 70
            level = "Easy"
        elif avg_words_per_sentence < 20:
            score = 50
            level = "Medium"
        elif avg_words_per_sentence < 25:
            score = 30
            level = "Difficult"
        else:
            score = 10
            level = "Very Difficult"
            
        return {"score": score, "level": level}
    
    def _check_grammar(self):
        """Basic grammar checking"""
        content = self.current_document["content"].lower()
        issues = []
        
        common_errors = [
            (r'\b(i)\s+(am|is|are|was|were)\s+(\w+ed)\b', "Passive voice detected"),
            (r'\b(very|really|quite|extremely)\s+(\w+)\b', "Consider removing intensifier"),
            (r'\b(there is|there are)\s+(\w+)\s+that\b', "Wordy construction"),
        ]
        
        for pattern, message in common_errors:
            if re.search(pattern, content):
                issues.append(message)
                
        return issues[:5]
    
    def _analyze_style(self):
        """Analyze writing style"""
        content = self.current_document["content"]
        
        style = "general"
        
        if re.search(r'\b(according to|citation|reference|study|research)\b', content, re.I):
            style = "academic"
        elif re.search(r'\b(proposal|budget|timeline|deliverable|stakeholder)\b', content, re.I):
            style = "business"
        elif re.search(r'\b(algorithm|function|class|import|def|return)\b', content):
            style = "technical"
        elif re.search(r'\b(chapter|scene|character|dialogue)\b', content, re.I):
            style = "creative"
            
        return {"detected_style": style, "confidence": 0.8}
    
    def _generate_suggestions(self):
        """Generate improvement suggestions"""
        content = self.current_document["content"]
        suggestions = []
        
        if len(content.split()) < 100:
            suggestions.append("Consider expanding the document with more details")
            
        structure = self._analyze_structure()
        if len(structure["headings"]) == 0 and len(content) > 500:
            suggestions.append("Add headings to improve document structure")
            
        readability = self._analyze_readability()
        if readability["score"] < 30:
            suggestions.append("Simplify sentences to improve readability")
            
        return suggestions


# ============================================================================
# MODEL ROUTER — centralised fallback for every LLM call in the app
# ============================================================================

# Ordered priority list.  Change this one list to update all components.
MODEL_PRIORITY: List[str] = [
    "llama-3.3-70b-versatile",
    "llama-3.1-70b-versatile",
    "mixtral-8x7b-32768",
]

class ModelRouter:
    """
    Drop-in helper that tries every model in MODEL_PRIORITY until one succeeds.
    Used by every intelligence-core component so there is a single, consistent
    fallback strategy across the whole app.

    Usage:
        router = ModelRouter(groq_client)

        # Non-streaming call — returns the text string
        text = router.complete(messages, max_tokens=500, temperature=0.2)

        # Streaming call — calls on_token for each chunk, returns full text
        text = router.stream(messages, on_token=callback, max_tokens=4000)

    Both methods return None (not raise) if every model fails, so callers can
    gracefully degrade to their regex / template fallback.
    """

    def __init__(self, client, models: Optional[List[str]] = None):
        self.client = client
        self.models = models or MODEL_PRIORITY
        # Expose which model actually answered (for debugging / UI display)
        self.last_model_used: Optional[str] = None
        self.last_attempt_errors: Dict[str, str] = {}

    def complete(
        self,
        messages: List[Dict],
        max_tokens: int = 1000,
        temperature: float = 0.3,
        timeout: int = 30,
    ) -> Optional[str]:
        """Try each model in order; return text on first success, None if all fail."""
        self.last_attempt_errors = {}
        for model in self.models:
            try:
                response = self.client.chat.completions.create(
                    model=model,
                    messages=messages,
                    max_tokens=max_tokens,
                    temperature=temperature,
                    timeout=timeout,
                )
                self.last_model_used = model
                return response.choices[0].message.content.strip()
            except Exception as e:
                self.last_attempt_errors[model] = str(e)
                print(f"[ModelRouter] {model} failed: {e}")
                continue
        print(f"[ModelRouter] All models failed: {self.last_attempt_errors}")
        return None

    def stream(
        self,
        messages: List[Dict],
        on_token: Optional[Callable[[str], None]] = None,
        max_tokens: int = 4000,
        temperature: float = 0.3,
        timeout: int = 60,
    ) -> Optional[str]:
        """
        Try streaming on each model.  If a model raises before/during streaming,
        fall through to the next.  Returns full accumulated text, or None.
        """
        self.last_attempt_errors = {}
        # Track speed metrics for caller
        self.last_stream_metrics: Dict[str, Any] = {}

        for model in self.models:
            full_content = ""
            token_count = 0
            first_token_time: Optional[float] = None
            start = time.time()
            try:
                stream = self.client.chat.completions.create(
                    model=model,
                    messages=messages,
                    max_tokens=max_tokens,
                    temperature=temperature,
                    stream=True,
                    timeout=timeout,
                )
                for chunk in stream:
                    if chunk.choices and chunk.choices[0].delta.content:
                        token = chunk.choices[0].delta.content
                        if first_token_time is None:
                            first_token_time = time.time()
                        full_content += token
                        token_count += 1
                        if on_token:
                            on_token(token)

                elapsed_ms = int((time.time() - start) * 1000)
                ttft_ms = int((first_token_time - start) * 1000) if first_token_time else 0
                tps = round(token_count / max(elapsed_ms / 1000, 0.001), 1)
                self.last_stream_metrics = {
                    "elapsed_ms": elapsed_ms,
                    "ttft_ms": ttft_ms,
                    "token_count": token_count,
                    "tokens_per_sec": tps,
                    "char_count": len(full_content),
                    "model": model,
                }
                self.last_model_used = model
                print(f"[ModelRouter] stream OK: {model} | {token_count} tok | {tps} tok/s")
                return full_content

            except Exception as e:
                self.last_attempt_errors[model] = str(e)
                print(f"[ModelRouter] stream {model} failed: {e}")
                # If we already got partial content and the stream broke mid-way,
                # treat it as a hard failure and try the next model cleanly.
                continue

        print(f"[ModelRouter] All stream models failed: {self.last_attempt_errors}")
        return None


# ============================================================================
# 1. CONVERSATION MANAGER
# ============================================================================

class ConversationManager:
    """Maintains persistent, multi-turn conversation memory with document state tracking"""
    
    def __init__(self, max_history: int = 10):
        self.max_history = max_history
        self.turns: List[ConversationTurn] = []
        self.intent_summary: Optional[str] = None
        self.cumulative_edits: Dict[str, Any] = {
            "total_edits": 0,
            "sections_affected": defaultdict(int),
            "first_interaction": None,
            "last_interaction": None
        }
    
    def add_user_message(self, query: str, document_state: dict) -> None:
        snapshot = {
            "word_count": document_state.get("word_count", 0),
            "char_count": document_state.get("char_count", 0),
            "title": document_state.get("title", "Untitled"),
            "content_preview": document_state.get("content", "")[:200],
            "has_content": bool(document_state.get("content", ""))
        }
        
        turn = ConversationTurn(
            user_query=query,
            document_snapshot=snapshot,
            assistant_response="",
            edits_made={}
        )
        self.turns.append(turn)
        
        if self.cumulative_edits["first_interaction"] is None:
            self.cumulative_edits["first_interaction"] = datetime.now()
        self.cumulative_edits["last_interaction"] = datetime.now()
        
        if len(self.turns) > self.max_history:
            self.turns = self.turns[-self.max_history:]
    
    def add_assistant_message(self, response: str, edits_made: dict) -> None:
        if self.turns:
            self.turns[-1].assistant_response = response
            self.turns[-1].edits_made = edits_made
            
            self.cumulative_edits["total_edits"] += 1
            for section in edits_made.get("sections_modified", []):
                self.cumulative_edits["sections_affected"][section] += 1
    
    def get_conversation_context(self) -> str:
        if not self.turns:
            return "No previous conversation."
        
        context_parts = ["## Conversation History\n"]
        
        for i, turn in enumerate(self.turns[-self.max_history:], 1):
            context_parts.append(f"**Turn {i}:**")
            context_parts.append(f"User: \"{turn.user_query[:200]}\"")
            context_parts.append(f"Document: {turn.document_snapshot.get('title', 'Untitled')} "
                               f"({turn.document_snapshot.get('word_count', 0)} words)")
            
            if turn.edits_made:
                changes_desc = ", ".join(turn.edits_made.get("key_changes", [])[:3])
                if changes_desc:
                    context_parts.append(f"Result: {changes_desc}")
            
            context_parts.append("")
        
        if self.cumulative_edits["total_edits"] > 1:
            context_parts.append(f"**Cumulative:** {self.cumulative_edits['total_edits']} edits across "
                               f"{len(self.cumulative_edits['sections_affected'])} sections")
        
        return "\n".join(context_parts)
    
    def get_document_evolution(self) -> List[Dict]:
        evolution = []
        for turn in self.turns:
            evolution.append({
                "turn_id": turn.turn_id,
                "timestamp": turn.timestamp.isoformat(),
                "query": turn.user_query[:100],
                "document_state": turn.document_snapshot,
                "changes": turn.edits_made
            })
        return evolution
    
    def summarize_intent(self, llm_client=None) -> str:
        if not self.turns:
            return "No conversation to summarize"
        
        queries = [turn.user_query for turn in self.turns[-5:]]
        
        if llm_client and len(queries) > 1:
            try:
                prompt = f"""Based on these user queries about document editing, what is the user's OVERARCHING intent?

Queries:
{chr(10).join(f'- {q}' for q in queries)}

Summarize in one sentence what the user is trying to achieve:"""
                messages = [{"role": "user", "content": prompt}]
                router = ModelRouter(llm_client)
                result = router.complete(messages, max_tokens=100, temperature=0.3)
                if result:
                    self.intent_summary = result
                    return self.intent_summary
            except Exception:
                pass
        
        keywords = []
        for q in queries:
            words = q.lower().split()[:5]
            keywords.extend(words)
        
        unique_keywords = list(set(keywords))[:5]
        self.intent_summary = f"User is focused on: {', '.join(unique_keywords)}"
        return self.intent_summary
    
    def clear(self) -> None:
        self.turns = []
        self.intent_summary = None
        self.cumulative_edits = {
            "total_edits": 0,
            "sections_affected": defaultdict(int),
            "first_interaction": None,
            "last_interaction": None
        }


# ============================================================================
# 2. INSTRUCTION PARSER
# ============================================================================

class InstructionParser:
    """Extract semantic meaning from user instructions using AI"""
    
    def __init__(self, llm_client=None):
        self.llm_client = llm_client
        self.confidence_threshold = 0.7
    
    def parse(self, instruction: str, document: dict) -> ParsedInstruction:
        if self.llm_client:
            try:
                return self._parse_with_ai(instruction, document)
            except Exception as e:
                print(f"AI parsing failed: {e}")
        
        return self._parse_with_regex(instruction, document)
    
    def _parse_with_ai(self, instruction: str, document: dict) -> ParsedInstruction:
        system_prompt = """You are an instruction parser for a document editing AI. Given a user instruction, extract semantic intent.

Output ONLY valid JSON with this structure:
{
    "intent": "improve|analyze|transform|generate|create|edit",
    "target_audience": "string or null",
    "tone": "formal|casual|academic|persuasive|neutral",
    "scope": "full|introduction|section|conclusion|paragraph",
    "domain": "business|academic|technical|creative|general",
    "constraints": ["list of specific requirements"],
    "confidence": 0.95,
    "reasoning": "Why this interpretation?",
    "needs_clarification": false,
    "clarification_questions": ["Question if needed?"],
    "extracted_entities": {"key": "value"}
}"""
        
        user_prompt = f"""Instruction: "{instruction}"

Document context: {document.get('title', 'Untitled')} ({document.get('word_count', 0)} words)

Parse this instruction and include clarification_questions if the instruction is ambiguous (confidence < 0.7):"""
        
        try:
            messages = [
                {"role": "system", "content": system_prompt},
                {"role": "user", "content": user_prompt}
            ]
            router = ModelRouter(self.llm_client)
            result_text = router.complete(messages, max_tokens=500, temperature=0.2)
            if result_text:
                json_match = re.search(r'\{.*\}', result_text, re.DOTALL)
                if json_match:
                    data = json.loads(json_match.group())
                    return ParsedInstruction(
                        intent=data.get("intent", "edit"),
                        target_audience=data.get("target_audience"),
                        tone=data.get("tone", "neutral"),
                        scope=data.get("scope", "full"),
                        domain=data.get("domain", "general"),
                        constraints=data.get("constraints", []),
                        reasoning=data.get("reasoning", ""),
                        confidence=data.get("confidence", 0.5),
                        needs_clarification=data.get("needs_clarification", False),
                        clarification_questions=data.get("clarification_questions", []),
                        extracted_entities=data.get("extracted_entities", {})
                    )
        except Exception as e:
            print(f"AI parsing error: {e}")
        
        return self._parse_with_regex(instruction, document)
    
    def _parse_with_regex(self, instruction: str, document: dict) -> ParsedInstruction:
        inst_lower = instruction.lower()
        
        intent = "edit"
        if any(word in inst_lower for word in ["improve", "enhance", "better"]):
            intent = "improve"
        elif any(word in inst_lower for word in ["analyze", "review", "check"]):
            intent = "analyze"
        elif any(word in inst_lower for word in ["transform", "convert", "change to"]):
            intent = "transform"
        elif any(word in inst_lower for word in ["generate", "create", "make"]):
            intent = "generate"
        
        tone = "neutral"
        if any(word in inst_lower for word in ["formal", "professional", "business"]):
            tone = "formal"
        elif any(word in inst_lower for word in ["casual", "friendly", "conversational"]):
            tone = "casual"
        elif any(word in inst_lower for word in ["academic", "scholarly", "research"]):
            tone = "academic"
        elif any(word in inst_lower for word in ["persuasive", "convincing", "compelling"]):
            tone = "persuasive"
        
        target_audience = None
        audience_patterns = [
            (r"for\s+a\s+(\d+[\s-]*year[\s-]*old)", "child"),
            (r"for\s+(executives|leaders|managers)", "executive"),
            (r"for\s+(beginners|novices)", "beginner"),
            (r"for\s+(experts|professionals)", "expert"),
        ]
        for pattern, audience_type in audience_patterns:
            match = re.search(pattern, inst_lower)
            if match:
                target_audience = match.group(1) if match.groups() else audience_type
                break
        
        scope = "full"
        if "introduction" in inst_lower:
            scope = "introduction"
        elif "conclusion" in inst_lower:
            scope = "conclusion"
        elif re.search(r'section\s+(\d+)', inst_lower):
            scope = f"section:{re.search(r'section\s+(\d+)', inst_lower).group(1)}"
        
        constraints = []
        word_match = re.search(r'under\s+(\d+)\s+words', inst_lower)
        if word_match:
            constraints.append(f"keep under {word_match.group(1)} words")
        
        if "keep accuracy" in inst_lower:
            constraints.append("preserve technical accuracy")
        
        confidence = 0.5
        if intent != "edit":
            confidence += 0.1
        if tone != "neutral":
            confidence += 0.1
        if constraints:
            confidence += 0.1
        confidence = min(confidence, 0.9)

        # FIX #9: Generate clarification questions when ambiguous
        clarification_questions = []
        needs_clarification = confidence < self.confidence_threshold

        if needs_clarification:
            if tone == "neutral":
                clarification_questions.append("What tone should I use? (formal, casual, academic, persuasive)")
            if target_audience is None:
                clarification_questions.append("Who is the target audience for this document?")
            if scope == "full" and len(instruction.split()) < 5:
                clarification_questions.append("Should I edit the full document or a specific section?")
            if not constraints:
                clarification_questions.append("Are there any length or style constraints I should follow?")
        
        return ParsedInstruction(
            intent=intent,
            target_audience=target_audience,
            tone=tone,
            scope=scope,
            domain="general",
            constraints=constraints,
            reasoning="Parsed using pattern matching",
            confidence=confidence,
            needs_clarification=needs_clarification,
            clarification_questions=clarification_questions[:3]
        )


# ============================================================================
# 3. FILE CONTEXT ACCUMULATOR
# ============================================================================

class FileContextAccumulator:
    """Remember all uploaded files and cross-reference them"""
    
    def __init__(self, llm_client=None):
        self.llm_client = llm_client
        self.files: Dict[str, Dict] = {}
        self.file_summaries: Dict[str, str] = {}
        self.semantic_index: Dict[str, List[str]] = defaultdict(list)
    
    def add_file(self, filename: str, file_type: str, content: str, metadata: dict) -> None:
        summary = self._generate_file_summary(filename, file_type, content, metadata)
        keywords = self._extract_keywords(content, metadata)
        
        self.files[filename] = {
            "filename": filename,
            "type": file_type,
            "content": content[:3000],
            "metadata": metadata,
            "summary": summary,
            "keywords": keywords,
            "timestamp": datetime.now().isoformat()
        }
        
        self.file_summaries[filename] = summary
        
        for keyword in keywords:
            self.semantic_index[keyword].append(filename)
    
    def _generate_file_summary(self, filename: str, file_type: str, content: str, metadata: dict) -> str:
        if file_type == "csv":
            lines = content.strip().split('\n')
            if len(lines) > 1:
                headers = lines[0].split(',')
                return f"CSV with {len(lines)-1} data rows, columns: {', '.join(headers[:5])}"
        elif file_type == "json":
            try:
                data = json.loads(content[:1000])
                if isinstance(data, dict):
                    return f"JSON object with keys: {', '.join(list(data.keys())[:5])}"
                elif isinstance(data, list):
                    return f"JSON array with {len(data)} items"
            except:
                pass
        
        words = len(content.split())
        return f"File with {words} words. Type: {file_type}"
    
    def _extract_keywords(self, content: str, metadata: dict) -> List[str]:
        keywords = set()
        if "columns" in metadata:
            keywords.update(metadata["columns"])
        
        words = content.lower().split()[:200]
        common_words = {"the", "a", "an", "and", "or", "but", "in", "on", "at", "to", "for"}
        
        for word in words:
            if len(word) > 3 and word not in common_words:
                keywords.add(word)
        
        return list(keywords)[:20]
    
    def get_file_context(self) -> str:
        if not self.files:
            return "No files uploaded."
        
        context_parts = ["## Uploaded Files Context\n"]
        for filename, file_info in self.files.items():
            context_parts.append(f"**File:** {filename}")
            context_parts.append(f"Type: {file_info['type']}")
            context_parts.append(f"Summary: {file_info['summary']}")
            context_parts.append("")
        
        return "\n".join(context_parts)
    
    def get_detailed_file_context(self, filename: str = None) -> str:
        if filename and filename in self.files:
            file_info = self.files[filename]
            return f"""## File: {filename}
Type: {file_info['type']}
Summary: {file_info['summary']}
Content Preview:
{file_info['content'][:500]}
"""
        
        result = ""
        for filename, file_info in self.files.items():
            result += f"\n### {filename}\n{file_info['summary']}\n"
        return result or "No files uploaded."
    
    def find_relevant_file(self, query: str) -> Optional[str]:
        query_lower = query.lower()
        best_match = None
        best_score = 0
        
        for filename, file_info in self.files.items():
            score = 0
            for keyword in file_info["keywords"]:
                if keyword in query_lower:
                    score += 1
            if filename.lower() in query_lower:
                score += 2
            if any(word in query_lower for word in file_info["summary"].lower().split()[:10]):
                score += 1
            
            if score > best_score and score > 0:
                best_score = score
                best_match = filename
        
        return best_match
    
    def get_cross_references(self, document_content: str) -> List[Dict]:
        suggestions = []
        for filename, file_info in self.files.items():
            doc_lower = document_content.lower()
            file_keywords = file_info["keywords"][:5]
            matched_keywords = [kw for kw in file_keywords if kw in doc_lower]
            if matched_keywords:
                suggestions.append({
                    "file": filename,
                    "type": file_info["type"],
                    "matched_terms": matched_keywords,
                    "suggestion": f"Reference data from {filename} regarding {', '.join(matched_keywords[:3])}"
                })
        return suggestions
    
    def clear(self) -> None:
        self.files.clear()
        self.file_summaries.clear()
        self.semantic_index.clear()


# ============================================================================
# 4. EDIT PLANNER
# ============================================================================

class EditPlanner:
    """Plan document transformations before executing them"""
    
    def __init__(self, llm_client=None):
        self.llm_client = llm_client
    
    def plan(self, parsed: ParsedInstruction, document: dict, conversation=None) -> EditPlan:
        if self.llm_client:
            try:
                return self._plan_with_ai(parsed, document, conversation)
            except Exception as e:
                print(f"AI planning failed: {e}")
        
        return self._plan_with_templates(parsed, document, conversation)
    
    def _plan_with_ai(self, parsed: ParsedInstruction, document: dict, conversation) -> EditPlan:
        conv_context = ""
        if conversation:
            conv_context = conversation.get_conversation_context()
        
        system_prompt = """You are an edit planner for a document AI. Create a detailed execution plan.

Output JSON:
{
    "strategy": "Overall approach description",
    "steps": ["Step 1", "Step 2", "Step 3"],
    "constraints": ["Constraint 1", "Constraint 2"],
    "target_metrics": {"metric": "value"},
    "rationale": "Why this approach"
}"""
        
        user_prompt = f"""Parsed Instruction:
- Intent: {parsed.intent}
- Audience: {parsed.target_audience}
- Tone: {parsed.tone}
- Scope: {parsed.scope}

Document: {document.get('title', 'Untitled')} ({document.get('word_count', 0)} words)

{conv_context}

Create edit plan:"""
        
        try:
            messages = [
                {"role": "system", "content": system_prompt},
                {"role": "user", "content": user_prompt}
            ]
            router = ModelRouter(self.llm_client)
            result_text = router.complete(messages, max_tokens=800, temperature=0.4)
            if result_text:
                json_match = re.search(r'\{.*\}', result_text, re.DOTALL)
                if json_match:
                    data = json.loads(json_match.group())
                    return EditPlan(
                        strategy=data.get("strategy", "Apply requested edits"),
                        steps=data.get("steps", ["Analyze document", "Apply changes", "Verify result"]),
                        constraints=data.get("constraints", parsed.constraints),
                        target_metrics=data.get("target_metrics", {}),
                        rationale=data.get("rationale", "Based on user instruction")
                    )
        except Exception as e:
            print(f"AI planning error: {e}")
        
        return self._plan_with_templates(parsed, document, conversation)
    
    def _plan_with_templates(self, parsed: ParsedInstruction, document: dict, conversation) -> EditPlan:
        intent_plans = {
            "improve": {
                "strategy": "Enhance document quality by improving clarity, flow, and engagement",
                "steps": ["Identify areas needing improvement", "Rewrite for better clarity", "Enhance vocabulary", "Ensure consistent tone"]
            },
            "analyze": {
                "strategy": "Perform comprehensive document analysis without modifying content",
                "steps": ["Analyze document structure", "Evaluate content quality", "Check for grammar issues", "Generate recommendations"]
            },
            "transform": {
                "strategy": "Transform document style and tone according to requirements",
                "steps": ["Understand target style", "Rewrite to match desired tone", "Adjust vocabulary", "Preserve core meaning"]
            },
            "generate": {
                "strategy": "Generate new content based on document context",
                "steps": ["Analyze existing content", "Identify gaps", "Generate relevant content", "Integrate smoothly"]
            }
        }
        
        plan_template = intent_plans.get(parsed.intent, intent_plans["improve"])
        
        steps = plan_template["steps"].copy()
        if parsed.tone != "neutral":
            steps.append(f"Adjust content to {parsed.tone} tone")
        if parsed.scope != "full":
            steps.insert(1, f"Focus exclusively on {parsed.scope} section")
        
        constraints = parsed.constraints.copy()
        if parsed.target_audience:
            constraints.append(f"Target audience: {parsed.target_audience}")
        
        return EditPlan(
            strategy=plan_template["strategy"],
            steps=steps,
            constraints=constraints,
            target_metrics={"preserve_facts": True},
            rationale=f"Template-based plan for {parsed.intent} operation"
        )


# ============================================================================
# 5. STREAMING RESPONSE HANDLER
# ============================================================================

class StreamingResponseHandler:
    """Stream responses token-by-token instead of blocking"""
    
    def __init__(self, client):
        self.client = client
        # Item 1: Expose last-run speed metrics for UI display
        self.last_metrics: Dict[str, Any] = {}

    def stream_completion(
        self, 
        messages: List[Dict], 
        on_token: Optional[Callable[[str], None]] = None,
        model: str = "llama-3.3-70b-versatile",
        temperature: float = 0.3,
        max_tokens: int = 4000
    ) -> str:
        full_content = ""
        token_count = 0
        start_time = time.time()
        first_token_time: Optional[float] = None
        
        try:
            stream = self.client.chat.completions.create(
                model=model,
                messages=messages,
                temperature=temperature,
                max_tokens=max_tokens,
                stream=True,
                timeout=60
            )
            
            for chunk in stream:
                if chunk.choices and chunk.choices[0].delta.content:
                    token = chunk.choices[0].delta.content
                    if first_token_time is None:
                        first_token_time = time.time()
                    full_content += token
                    token_count += 1
                    if on_token:
                        on_token(token)
            
            elapsed_ms = int((time.time() - start_time) * 1000)
            ttft_ms = int((first_token_time - start_time) * 1000) if first_token_time else 0
            tokens_per_sec = round(token_count / max(elapsed_ms / 1000, 0.001), 1)

            # Item 1: Store metrics for caller to display
            self.last_metrics = {
                "elapsed_ms": elapsed_ms,
                "ttft_ms": ttft_ms,
                "token_count": token_count,
                "tokens_per_sec": tokens_per_sec,
                "char_count": len(full_content),
            }
            print(f"Streaming: {token_count} tokens, {tokens_per_sec} tok/s, TTFT {ttft_ms}ms")
            return full_content
            
        except Exception as e:
            print(f"Streaming error: {e}")
            try:
                response = self.client.chat.completions.create(
                    model=model,
                    messages=messages,
                    temperature=temperature,
                    max_tokens=max_tokens,
                    stream=False
                )
                full_content = response.choices[0].message.content.strip()
                self.last_metrics = {"elapsed_ms": 0, "ttft_ms": 0, "token_count": 0,
                                     "tokens_per_sec": 0, "char_count": len(full_content)}
                if on_token:
                    on_token(full_content)
                return full_content
            except Exception as e2:
                print(f"Fallback failed: {e2}")
                return f"Error: {str(e)}"


# ============================================================================
# 6. DOCUMENT PROFILER
# ============================================================================

class DocumentProfiler:
    """Deep AI-driven analysis of document (replaces regex analysis)"""
    
    def __init__(self, llm_client=None):
        self.llm_client = llm_client
    
    def profile(self, document: dict) -> DocumentProfile:
        content = document.get("content", "")
        
        if not content or len(content.strip()) < 50:
            return self._empty_profile("Document too short for analysis")
        
        if self.llm_client:
            try:
                return self._profile_with_ai(document)
            except Exception as e:
                print(f"AI profiling failed: {e}")
        
        return self._profile_with_fallback(document)
    
    def _profile_with_ai(self, document: dict) -> DocumentProfile:
        content = document.get("content", "")
        title = document.get("title", "Untitled")
        
        system_prompt = """You are a document profiler. Analyze the document and output JSON.

Output format:
{
    "structure": {
        "has_clear_intro": true/false,
        "has_body_paragraphs": true/false,
        "has_conclusion": true/false,
        "logical_flow": "good|fair|poor",
        "issues": ["specific structural issues"]
    },
    "content": {
        "primary_purpose": "inform|persuade|entertain|instruct",
        "target_audience": "inferred audience description",
        "tone": "formal|casual|academic|persuasive",
        "reading_level": "elementary|high_school|college|expert"
    },
    "quality": {
        "grammar_issues": ["specific grammar issues"],
        "clarity_problems": ["unclear sections"],
        "engagement_score": 0.0-1.0
    },
    "suggestions": ["specific, actionable suggestion 1", "suggestion 2"],
    "strengths": ["strength 1", "strength 2"]
}"""
        
        user_prompt = f"""Title: {title}

Content:
{content[:3000]}

Analyze this document:"""
        
        try:
            messages = [
                {"role": "system", "content": system_prompt},
                {"role": "user", "content": user_prompt}
            ]
            router = ModelRouter(self.llm_client)
            result_text = router.complete(messages, max_tokens=1000, temperature=0.3)
            if result_text:
                json_match = re.search(r'\{.*\}', result_text, re.DOTALL)
                if json_match:
                    data = json.loads(json_match.group())
                    return DocumentProfile(
                        structure=data.get("structure", {}),
                        content=data.get("content", {}),
                        quality=data.get("quality", {}),
                        suggestions=data.get("suggestions", []),
                        strengths=data.get("strengths", []),
                        metadata={"analyzed_by": f"AI ({router.last_model_used})",
                                  "timestamp": datetime.now().isoformat()}
                    )
        except Exception as e:
            print(f"AI profile error: {e}")
        
        return self._profile_with_fallback(document)
    
    def _profile_with_fallback(self, document: dict) -> DocumentProfile:
        content = document.get("content", "")
        
        has_intro = any(word in content[:500].lower() for word in ["introduction", "overview"])
        has_conclusion = any(word in content[-500:].lower() for word in ["conclusion", "summary"])
        paragraphs = [p for p in content.split('\n\n') if len(p.strip()) > 50]
        has_body = len(paragraphs) >= 2
        
        grammar_issues = []
        passive_matches = re.findall(r'\b(?:is|are|was|were|be|been|being)\s+\w+ed\b', content, re.I)
        if len(passive_matches) > 5:
            grammar_issues.append("Excessive passive voice usage")
        
        long_sentences = [s for s in re.split(r'[.!?]+', content) if len(s.split()) > 30]
        if len(long_sentences) > 3:
            grammar_issues.append(f"{len(long_sentences)} very long sentences")
        
        suggestions = []
        if not has_intro:
            suggestions.append("Add a clear introduction")
        if not has_conclusion:
            suggestions.append("Add a conclusion to summarize key points")
        if len(content.split()) < 200:
            suggestions.append("Consider expanding with more details")
        
        tone = "neutral"
        formal_count = sum(content.lower().count(w) for w in ["therefore", "consequently"])
        casual_count = sum(content.lower().count(w) for w in ["basically", "actually"])
        if formal_count > casual_count * 2:
            tone = "formal"
        elif casual_count > formal_count * 2:
            tone = "casual"
        
        return DocumentProfile(
            structure={
                "has_clear_intro": has_intro,
                "has_body_paragraphs": has_body,
                "has_conclusion": has_conclusion,
                "logical_flow": "fair" if has_intro and has_body else "poor",
                "issues": []
            },
            content={
                "primary_purpose": "inform",
                "target_audience": "general audience",
                "tone": tone,
                "reading_level": "college" if len(content.split()) > 500 else "high_school"
            },
            quality={
                "grammar_issues": grammar_issues[:3],
                "clarity_problems": [],
                "engagement_score": 0.5
            },
            suggestions=suggestions[:5],
            strengths=[],
            metadata={"analyzed_by": "fallback"}
        )
    
    def _empty_profile(self, reason: str) -> DocumentProfile:
        return DocumentProfile(
            structure={"has_clear_intro": False, "has_conclusion": False, "logical_flow": "poor", "issues": [reason]},
            content={"primary_purpose": "unknown", "target_audience": "unknown", "tone": "neutral", "reading_level": "unknown"},
            quality={"grammar_issues": [], "clarity_problems": [], "engagement_score": 0.0},
            suggestions=["Add more content for proper analysis"],
            strengths=[],
            metadata={"error": reason}
        )


# ============================================================================
# 7. CONTEXTUAL EDITOR (Core Edit Engine)
# ============================================================================

class ContextualEditor:
    """Execute document edits using full context - orchestrates all components"""
    
    def __init__(self, llm_client, streaming_handler: StreamingResponseHandler):
        self.llm_client = llm_client
        self.streaming_handler = streaming_handler
        self.instruction_parser = InstructionParser(llm_client)
        self.edit_planner = EditPlanner(llm_client)
    
    def edit(
        self,
        instruction: str,
        document: dict,
        conversation: Optional[ConversationManager] = None,
        file_context: Optional[FileContextAccumulator] = None,
        stream_callback: Optional[Callable[[str], None]] = None,
        selected_files: Optional[List[str]] = None,  # Item 3: multi-file selection
    ) -> EditResult:
        start_time = time.time()
        
        parsed = self.instruction_parser.parse(instruction, document)
        
        if parsed.needs_clarification and parsed.confidence < 0.6:
            return EditResult(
                edited_document=document.get("content", ""),
                changes_made={},
                reasoning=f"Need clarification",
                successful=False,
                execution_time_ms=int((time.time() - start_time) * 1000)
            )
        
        conv_context = ""
        if conversation:
            conv_context = conversation.get_conversation_context()
        
        file_context_str = ""
        if file_context and file_context.files:
            # Item 2: Auto-reference files — inject content of matched/selected files
            if selected_files:
                # User explicitly chose files (Item 3)
                for fname in selected_files:
                    detail = file_context.get_detailed_file_context(fname)
                    file_context_str += detail + "\n"
            else:
                # Auto-detect the most relevant file
                relevant_file = file_context.find_relevant_file(instruction)
                if relevant_file:
                    # Inject summary + content preview so AI can actually use the data
                    file_context_str = file_context.get_detailed_file_context(relevant_file)
                else:
                    file_context_str = file_context.get_file_context()
        
        plan = self.edit_planner.plan(parsed, document, conversation)
        
        prompt = self._build_edit_prompt(
            instruction=instruction,
            parsed=parsed,
            plan=plan,
            document=document,
            conv_context=conv_context,
            file_context=file_context_str
        )
        
        messages = [
            {"role": "system", "content": self._get_system_prompt(conversation)},
            {"role": "user", "content": prompt}
        ]
        
        router = ModelRouter(self.llm_client)

        try:
            if stream_callback:
                edited_content = router.stream(
                    messages=messages,
                    on_token=stream_callback,
                    max_tokens=4000,
                    temperature=0.3,
                )
                # Sync metrics to streaming_handler so UI can read them
                if router.last_stream_metrics:
                    self.streaming_handler.last_metrics = router.last_stream_metrics
            else:
                edited_content = router.complete(
                    messages=messages,
                    max_tokens=4000,
                    temperature=0.3,
                )

            if not edited_content:
                raise RuntimeError(
                    f"All models failed. Errors: {router.last_attempt_errors}"
                )

            changes_made = self._calculate_changes(
                document.get("content", ""),
                edited_content,
                plan
            )
            execution_ms = int((time.time() - start_time) * 1000)
            return EditResult(
                edited_document=edited_content,
                changes_made=changes_made,
                reasoning=plan.rationale,
                successful=True,
                execution_time_ms=execution_ms
            )

        except Exception as e:
            return EditResult(
                edited_document=document.get("content", ""),
                changes_made={},
                reasoning=f"Edit failed (all models tried): {str(e)}",
                successful=False,
                execution_time_ms=int((time.time() - start_time) * 1000)
            )
    
    def _build_edit_prompt(self, instruction, parsed, plan, document, conv_context, file_context):
        content_preview = document.get("content", "")
        if len(content_preview) > 4000:
            content_preview = content_preview[:4000] + "\n...[truncated]..."
        
        prompt_parts = [
            "## EDIT INSTRUCTION",
            f"User: {instruction}",
            "",
            "## PARSED INTENT",
            f"- Intent: {parsed.intent}",
            f"- Tone: {parsed.tone}",
            f"- Audience: {parsed.target_audience or 'Not specified'}",
            f"- Scope: {parsed.scope}",
            f"- Constraints: {', '.join(parsed.constraints) if parsed.constraints else 'None'}",
            "",
            "## EDIT PLAN",
            f"Strategy: {plan.strategy}",
            f"Steps:",
        ]
        
        for step in plan.steps:
            prompt_parts.append(f"  {step}")

        # FIX #5: Make conversation context PROMINENT at top of prompt
        if conv_context and conv_context != "No previous conversation.":
            prompt_parts.extend([
                "",
                "## ⚠️ CRITICAL - CONVERSATION CONTEXT (MUST FOLLOW)",
                conv_context,
                "IMPORTANT: The above history shows what was done previously.",
                "You MUST maintain any tone/style/constraints established in previous turns.",
            ])
        
        if file_context and file_context != "No files uploaded.":
            prompt_parts.extend(["", file_context])
        
        prompt_parts.extend([
            "",
            "## DOCUMENT TO EDIT",
            "```",
            content_preview,
            "```",
            "",
            "Return ONLY the edited document content."
        ])
        
        return "\n".join(prompt_parts)
    
    def _get_system_prompt(self, conversation=None) -> str:
        # FIX #5: Inject cumulative intent into system prompt
        intent_note = ""
        if conversation and conversation.intent_summary:
            intent_note = f"\nUser's overarching goal: {conversation.intent_summary}\nMaintain this goal across all edits."

        return f"""You are MozeAI Document Editor, a precise document editing AI.{intent_note}

Return ONLY the edited document content - no explanations, no chat responses.
Preserve the original meaning unless instructed otherwise.
Apply changes exactly as described.
ALWAYS maintain any tone, style, or constraints established in previous conversation turns."""
    
    def _calculate_changes(self, old_content: str, new_content: str, plan: EditPlan) -> Dict:
        old_words = len(old_content.split())
        new_words = len(new_content.split())
        word_diff = new_words - old_words
        
        return {
            "additions": max(0, word_diff),
            "deletions": max(0, -word_diff),
            "net_change": word_diff,
            "old_word_count": old_words,
            "new_word_count": new_words,
            "sections_affected": ["content"],
            "key_changes": [f"Word count: {word_diff:+d} words ({old_words} → {new_words})"]
        }


# ============================================================================
# FILE PROCESSING FUNCTIONS
# ============================================================================

def extract_text_from_pdf(file):
    try:
        file.seek(0)
        pdf_reader = PyPDF2.PdfReader(file)
        text = ""
        for page_num, page in enumerate(pdf_reader.pages):
            page_text = page.extract_text()
            if page_text and page_text.strip():
                text += f"\n--- Page {page_num + 1} ---\n"
                text += page_text.strip() + "\n"
        return text[:5000] if text.strip() else "No extractable text in PDF"
    except Exception as e:
        return f"Error reading PDF: {str(e)}"

def extract_text_from_docx(file):
    try:
        file.seek(0)
        doc = docx.Document(file)
        text = ""
        for para in doc.paragraphs:
            if para.text and para.text.strip():
                text += para.text.strip() + "\n\n"
        return text[:5000] if text.strip() else "No extractable text in document"
    except Exception as e:
        return f"Error reading Word document: {str(e)}"

def extract_text_from_txt(file):
    try:
        file.seek(0)
        content = file.read().decode('utf-8')
        return content[:5000] if content.strip() else "File is empty"
    except UnicodeDecodeError:
        try:
            file.seek(0)
            content = file.read().decode('latin-1')
            return content[:5000]
        except:
            return "Error decoding text file"
    except Exception as e:
        return f"Error reading text file: {str(e)}"

def extract_text_from_csv(file):
    try:
        file.seek(0)
        content = file.read().decode('utf-8')
        csv_reader = csv.reader(StringIO(content))
        text = "CSV Data:\n\n"
        rows = list(csv_reader)
        if rows:
            text += "Headers: " + " | ".join(rows[0]) + "\n\n"
            for i, row in enumerate(rows[1:11], 1):
                text += f"Row {i}: " + " | ".join(row) + "\n"
        return text[:5000] if text.strip() else "CSV file appears empty"
    except Exception as e:
        return f"Error reading CSV: {str(e)}"

def extract_text_from_json(file):
    try:
        file.seek(0)
        content = file.read().decode('utf-8')
        data = json.loads(content)
        formatted = json.dumps(data, indent=2)
        return formatted[:5000] if formatted else "JSON file is empty"
    except Exception as e:
        return f"Error reading JSON: {str(e)}"

def process_uploaded_file(uploaded_file):
    file_type = uploaded_file.type
    file_name = uploaded_file.name.lower()
    
    if file_type == "application/pdf" or file_name.endswith('.pdf'):
        return extract_text_from_pdf(uploaded_file)
    elif file_type == "application/vnd.openxmlformats-officedocument.wordprocessingml.document" or file_name.endswith('.docx'):
        return extract_text_from_docx(uploaded_file)
    elif file_type == "text/plain" or file_name.endswith('.txt'):
        return extract_text_from_txt(uploaded_file)
    elif file_type == "text/csv" or file_name.endswith('.csv'):
        return extract_text_from_csv(uploaded_file)
    elif file_type == "application/json" or file_name.endswith('.json'):
        return extract_text_from_json(uploaded_file)
    else:
        return f"Unsupported file type: {file_type}"

# ============================================================================
# DOCUMENT GENERATION FUNCTIONS
# ============================================================================

def create_ppt_from_content(title, content, filename="presentation"):
    try:
        prs = Presentation()
        title_slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_slide_layout)
        slide.shapes.title.text = title[:100]
        slide.placeholders[1].text = f"Created by MozeAI\n{datetime.now().strftime('%Y-%m-%d %H:%M:%S')}"
        
        content_slide_layout = prs.slide_layouts[1]
        lines = content.split('\n')
        current_slide = None
        current_text_frame = None
        
        for line in lines:
            line = line.strip()
            if not line:
                continue
            
            if len(line) < 60 and (line.endswith(':') or line.isupper() or re.match(r'^\d+\.', line)):
                current_slide = prs.slides.add_slide(content_slide_layout)
                current_slide.shapes.title.text = line.rstrip(':')[:100]
                content_box = current_slide.placeholders[1]
                current_text_frame = content_box.text_frame
                current_text_frame.text = ""
            else:
                if current_slide is None:
                    current_slide = prs.slides.add_slide(content_slide_layout)
                    current_slide.shapes.title.text = "Content"
                    content_box = current_slide.placeholders[1]
                    current_text_frame = content_box.text_frame
                    current_text_frame.text = ""
                
                if current_text_frame:
                    p = current_text_frame.add_paragraph()
                    p.text = line[:150]
                    p.font.size = Pt(18)
        
        ppt_bytes = BytesIO()
        prs.save(ppt_bytes)
        ppt_bytes.seek(0)
        return ppt_bytes
    except Exception as e:
        print(f"PPT error: {e}")
        return None

def create_word_from_content(title, content, filename="document"):
    try:
        doc = WordDocument()
        title_heading = doc.add_heading(title, 0)
        title_heading.alignment = WD_ALIGN_PARAGRAPH.CENTER
        doc.add_paragraph(f"Generated by MozeAI on {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}")
        doc.add_paragraph()
        
        paragraphs = content.split('\n\n')
        for para in paragraphs:
            if para.strip():
                doc.add_paragraph(para.strip())
        
        word_bytes = BytesIO()
        doc.save(word_bytes)
        word_bytes.seek(0)
        return word_bytes
    except Exception as e:
        return None

def create_real_excel_file(title, data_rows):
    try:
        from openpyxl import Workbook
        from openpyxl.styles import Font, PatternFill, Alignment
        from openpyxl.utils import get_column_letter
        
        wb = Workbook()
        ws = wb.active
        ws.title = title[:31].replace('/', '_')
        
        for row_idx, row in enumerate(data_rows, 1):
            for col_idx, value in enumerate(row, 1):
                cell = ws.cell(row=row_idx, column=col_idx, value=value)
                if row_idx == 1:
                    cell.font = Font(bold=True, color="FFFFFF")
                    cell.fill = PatternFill(start_color="4472C4", end_color="4472C4", fill_type="solid")
        
        for col in ws.columns:
            max_length = 0
            for cell in col:
                try:
                    if len(str(cell.value)) > max_length:
                        max_length = len(str(cell.value))
                except:
                    pass
            ws.column_dimensions[get_column_letter(col[0].column)].width = min(max_length + 2, 50)
        
        output = BytesIO()
        wb.save(output)
        output.seek(0)
        return output
    except Exception as e:
        print(f"Excel error: {e}")
        return None

def create_csv_from_data(title, data_rows):
    try:
        output = BytesIO()
        output.write('\ufeff'.encode('utf-8'))
        writer = csv.writer(output)
        for row in data_rows:
            writer.writerow(row)
        output.seek(0)
        return output
    except Exception as e:
        return None

def export_chat_history():
    """FIX #10: Enhanced export includes intent + timeline"""
    if not st.session_state.chat_history:
        return None
    
    export_content = "=" * 70 + "\n"
    export_content += "CHAT HISTORY WITH MOZEAI\n"
    export_content += f"Exported on: {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}\n"
    export_content += "=" * 70 + "\n\n"
    
    # Include intent summary
    core = st.session_state.get("intelligence_core")
    if core:
        conv_manager = core["conversation_manager"]
        if conv_manager.intent_summary:
            export_content += f"=== SESSION INTENT ===\n{conv_manager.intent_summary}\n\n"
        
        evolution = conv_manager.get_document_evolution()
        if evolution:
            export_content += "=== EDIT TIMELINE ===\n"
            for i, turn in enumerate(evolution, 1):
                changes = turn.get("changes", {})
                net = changes.get("net_change", 0)
                export_content += f"{i}. {turn['query'][:60]}... → {net:+d} words\n"
            export_content += "\n"
    
    export_content += "=== CONVERSATION ===\n\n"
    for idx, (role, msg) in enumerate(st.session_state.chat_history, 1):
        if role == "user":
            export_content += f"[{idx}] USER:\n{msg}\n\n"
        else:
            export_content += f"[{idx}] MOZEAI:\n{msg}\n\n"
    
    return export_content

# ============================================================================
# WEB & UTILITY FUNCTIONS
# ============================================================================

def get_current_datetime():
    tz = pytz.timezone('Asia/Seoul')
    now = datetime.now(tz)
    return f"Date: {now.strftime('%B %d, %Y')}\nTime: {now.strftime('%I:%M %p')}\nTimezone: Asia/Seoul"

def internet_search(query):
    try:
        clean_query = query.strip()
        url = "https://html.duckduckgo.com/html/"
        params = {"q": clean_query}
        headers = {"User-Agent": "Mozilla/5.0"}
        response = requests.post(url, data=params, headers=headers, timeout=10)
        
        if response.status_code == 200:
            results = re.findall(r'<a rel="nofollow" class="result__a" href="[^"]*">([^<]+)</a>', response.text)
            snippets = re.findall(r'<a class="result__snippet"[^>]*>([^<]+)</a>', response.text)
            
            if results:
                context = f"SEARCH RESULTS for '{clean_query}':\n\n"
                for i in range(min(3, len(results))):
                    context += f"- {results[i]}\n"
                    if i < len(snippets):
                        snippet = re.sub(r'<[^>]+>', '', snippets[i])
                        context += f"  {snippet[:300]}...\n\n"
                return context[:2000]
        return ""
    except:
        return ""

def generate_image_with_quality(prompt, quality="high", style="realistic"):
    try:
        enhanced_prompt = f"{prompt}, high quality, detailed"
        encoded_prompt = requests.utils.quote(enhanced_prompt)
        timestamp = int(time.time())
        image_url = f"https://image.pollinations.ai/prompt/{encoded_prompt}?width=1024&height=1024&seed={timestamp}"
        return image_url
    except Exception as e:
        return None

def generate_and_display_image(prompt, is_edit=False):
    image_url = generate_image_with_quality(prompt)
    if image_url:
        return f"Generated Image for: '{prompt}'\n\n![Image]({image_url})"
    else:
        return "Sorry, I couldn't generate an image."

def llm_with_fallback(messages, max_retries=2):
    """General-purpose fallback wrapper — now delegates to ModelRouter for consistency."""
    router = ModelRouter(client)
    result = router.complete(messages, max_tokens=800, temperature=0.3, timeout=30)
    if result:
        st.session_state.last_model_used = router.last_model_used
        return result
    return "AI service temporarily unavailable."

def reason(question, context):
    messages = [
        {"role": "system", "content": "You are MozeAI, a helpful AI assistant."},
        {"role": "user", "content": f"{context}\n\nUSER QUESTION: {question}\n\nANSWER:"}
    ]
    return llm_with_fallback(messages)

def clean_answer(text):
    text = text.split("🧠")[0]
    text = text.split("Plan:")[0]
    return text.strip()


# ============================================================================
# ENHANCED RUN AGENT WITH INTELLIGENCE CORE
# ============================================================================

def handle_intelligent_edit(instruction: str, stream_callback=None, selected_files=None):
    """Use contextual editor for intelligent document editing"""
    core = st.session_state.intelligence_core
    
    document_state = {
        "content": st.session_state.workspace.current_document["content"],
        "title": st.session_state.workspace.current_document["title"],
        "word_count": len(st.session_state.workspace.current_document["content"].split()),
        "char_count": len(st.session_state.workspace.current_document["content"])
    }

    # Check clarification BEFORE editing
    parser = core["instruction_parser"]
    parsed = parser.parse(instruction, document_state)

    if parsed.needs_clarification and parsed.confidence < 0.6 and parsed.clarification_questions:
        st.session_state.pending_clarification = {
            "instruction": instruction,
            "questions": parsed.clarification_questions,
            "parsed": parsed
        }
        return None  # Signal that clarification is needed
    
    result = core["contextual_editor"].edit(
        instruction=instruction,
        document=document_state,
        conversation=core["conversation_manager"],
        file_context=core["file_accumulator"],
        stream_callback=stream_callback,
        selected_files=selected_files,  # Item 2 & 3: pass through
    )
    
    if result.successful:
        core["conversation_manager"].add_user_message(instruction, document_state)
        core["conversation_manager"].add_assistant_message("Document edited", result.changes_made)
        st.session_state.workspace.update_document(result.edited_document, f"AI Edit: {instruction[:100]}")
        
    return result

def profile_current_document():
    """Profile current document using AI"""
    core = st.session_state.intelligence_core
    document_state = {
        "content": st.session_state.workspace.current_document["content"],
        "title": st.session_state.workspace.current_document["title"]
    }
    return core["document_profiler"].profile(document_state)

def run_agent(query: str, stream_callback=None, selected_files=None):
    """Enhanced agent with intelligence core integration"""
    q = query.lower().strip()
    
    # Document workspace commands
    if q.startswith("/"):
        return handle_document_command(q)
    
    # Check for document editing commands
    edit_keywords = ["improve", "rewrite", "summarize", "expand", "shorten", 
                     "fix grammar", "make formal", "make academic", "translate"]
    
    if any(keyword in q for keyword in edit_keywords) and st.session_state.workspace.current_document["content"]:
        # Item 6: Edit validation — snapshot word count before edit
        pre_word_count = len(st.session_state.workspace.current_document["content"].split())

        result = handle_intelligent_edit(query, stream_callback=stream_callback, selected_files=selected_files)

        # Clarification needed
        if result is None:
            return "__CLARIFICATION_NEEDED__"

        if result.successful:
            # Item 6: Validate edit didn't produce empty/trivially-short output
            post_word_count = len(result.edited_document.split())
            if post_word_count < max(10, pre_word_count * 0.1):
                result.successful = False
                result.reasoning = (
                    f"Edit validation failed: output was only {post_word_count} words "
                    f"(original was {pre_word_count}). Original preserved."
                )
                st.session_state.workspace.update_document(
                    st.session_state.workspace.version_history[-1]["content"],
                    "Rollback: edit validation failed"
                )
                return result

            # Compute file cross-references
            doc_content = st.session_state.workspace.current_document["content"]
            file_acc = st.session_state.intelligence_core["file_accumulator"]
            cross_refs = file_acc.get_cross_references(doc_content)
            st.session_state.pending_file_suggestions = cross_refs

            return result
        else:
            return f"⚠️ {result.reasoning}"
    
    # Analysis command
    if "analyze document" in q or "profile document" in q:
        with st.spinner("Analyzing document..."):
            profile = profile_current_document()
            st.session_state.doc_profile_cache = profile
            
            result = f"## Document Analysis\n\n"
            result += f"**Tone:** {profile.content.get('tone', 'unknown').title()}\n"
            result += f"**Purpose:** {profile.content.get('primary_purpose', 'unknown').title()}\n"
            result += f"**Reading Level:** {profile.content.get('reading_level', 'unknown')}\n\n"
            
            if profile.strengths:
                result += "**Strengths:**\n"
                for s in profile.strengths[:3]:
                    result += f"- {s}\n"
                result += "\n"
            
            if profile.suggestions:
                result += "**Suggestions:**\n"
                for s in profile.suggestions[:3]:
                    result += f"- {s}\n"
            
            return result
    
    # Clear context
    if any(phrase in q for phrase in ["clear context", "new chat", "start fresh"]):
        st.session_state.workspace = DocumentWorkspace()
        st.session_state.intelligence_core["conversation_manager"].clear()
        st.session_state.intelligence_core["file_accumulator"].clear()
        st.session_state.chat_history = []
        st.session_state.uploaded_files = {}
        st.session_state.pending_clarification = None
        st.session_state.pending_file_suggestions = []
        return "✨ Everything cleared! Ready for a new session."
    
    # What is a word?
    if q == "what is a word":
        return "A **Word document** (.docx) is created by Microsoft Word. Try 'make a word about dogs'"
    
    # Excel generation
    if "make an excel" in q or "create an excel" in q or "generate an excel" in q:
        topic = q.replace("make an excel", "").replace("create an excel", "").replace("generate an excel", "").strip()
        topic = topic or "Sample_Data"
        
        data_rows = [
            ["Item", "Category", "Quantity", "Price", "Total"],
            ["Product A", "Electronics", 10, 99.99, 999.90],
            ["Product B", "Clothing", 25, 49.99, 1249.75],
            ["Product C", "Food", 50, 9.99, 499.50]
        ]
        
        excel_data = create_real_excel_file(topic, data_rows)
        if excel_data:
            st.session_state.excel_data = excel_data
            st.session_state.excel_topic = topic
            st.session_state.show_excel_download = True
            return f"📊 Created Excel file: {topic}. Scroll down to download!"
    
    # PowerPoint generation
    if any(phrase in q for phrase in ["make a ppt", "create a powerpoint"]):
        topic = q.replace("make a ppt", "").replace("create a powerpoint", "").strip() or "Presentation"
        content = f"Introduction to {topic}\n- Key point 1\n- Key point 2\n\nConclusion\n- Summary"
        ppt_bytes = create_ppt_from_content(topic, content)
        if ppt_bytes:
            st.session_state.ppt_data = ppt_bytes
            st.session_state.ppt_topic = topic
            st.session_state.show_ppt_download = True
            return f"📊 Created PowerPoint: {topic}. Scroll down to download!"
    
    # Word generation
    if any(phrase in q for phrase in ["make a word", "create a document"]):
        topic = q.replace("make a word", "").replace("create a document", "").strip() or "Document"
        content = f"# {topic}\n\nThis document covers important information about {topic}.\n\n## Introduction\n\nContent here.\n\n## Conclusion\n\nSummary."
        word_bytes = create_word_from_content(topic, content)
        if word_bytes:
            st.session_state.word_data = word_bytes
            st.session_state.word_topic = topic
            st.session_state.show_word_download = True
            return f"📄 Created Word document: {topic}. Scroll down to download!"
    
    # Image generation
    if any(phrase in q for phrase in ["generate image", "create image"]):
        image_prompt = q.replace("generate image", "").replace("create image", "").strip()
        if not image_prompt:
            image_prompt = "a beautiful landscape"
        return generate_and_display_image(image_prompt)
    
    # Item 9: Streaming for non-edit queries — use ModelRouter so fallback applies here too
    search_result = internet_search(query)
    context = get_current_datetime()
    if search_result:
        context += "\n" + search_result

    messages = [
        {"role": "system", "content": "You are MozeAI, a helpful AI assistant."},
        {"role": "user", "content": f"{context}\n\nUSER QUESTION: {query}\n\nANSWER:"}
    ]

    if stream_callback:
        router = ModelRouter(client)
        result = router.stream(messages, on_token=stream_callback, max_tokens=800)
        if result:
            st.session_state.last_model_used = router.last_model_used
            return result
        return "AI service temporarily unavailable."
    else:
        return reason(query, context)

def handle_document_command(command: str) -> str:
    """Handle slash commands"""
    cmd = command.lower().strip()
    workspace = st.session_state.workspace
    core = st.session_state.intelligence_core
    
    if cmd == "/analyze":
        analysis = workspace.analyze_document()
        return f"""## Document Analysis

**Structure:** {len(analysis['structure']['headings'])} headings, {analysis['structure']['paragraph_count']} paragraphs
**Readability:** {analysis['readability']['level']}
**Style:** {analysis['style_analysis']['detected_style']}

**Suggestions:**
{chr(10).join(f'- {s}' for s in analysis['suggestions'])}"""
    
    elif cmd == "/stats":
        meta = workspace.current_document["metadata"]
        return f"""## Document Stats

**Title:** {workspace.current_document['title']}
**Words:** {meta['word_count']}
**Characters:** {meta['char_count']}
**Reading Time:** {meta['reading_time']} min
**Versions:** {len(workspace.version_history)}"""
    
    elif cmd.startswith("/version"):
        parts = cmd.split()
        if len(parts) > 1 and parts[1].isdigit():
            if workspace.restore_version(int(parts[1])):
                return f"✅ Restored version {parts[1]}"
        return f"Versions: {len(workspace.version_history)} saved"
    
    elif cmd == "/conversation":
        summary = core["conversation_manager"].summarize_intent(client)
        return f"**Conversation Intent:** {summary}\n**Turns:** {len(core['conversation_manager'].turns)}"
    
    elif cmd == "/help":
        return """## Commands

**Document:** `/analyze`, `/stats`, `/version N`
**Conversation:** `/conversation`, `/clear`
**Editing:** Just tell me what to do, like "make this formal" or "add a conclusion" """
    
    else:
        return f"Unknown command. Type `/help` for available commands."


# ============================================================================
# FIX #2: CLARIFICATION DIALOG COMPONENT
# ============================================================================

# UI written below

# ============================================================================
# CLARIFICATION DIALOG
# ============================================================================

def render_clarification_dialog():
    pending = st.session_state.get("pending_clarification")
    if not pending:
        return
    st.markdown(
        '<div class="clarif-box">'
        f'<div class="clarif-title">I need a bit more detail</div>'
        f'<div class="clarif-sub">Instruction: <em>{pending["instruction"]}</em></div>'
        '</div>',
        unsafe_allow_html=True,
    )
    answers = {}
    for i, question in enumerate(pending["questions"]):
        answer = st.text_input(question, key=f"clarif_q_{i}", placeholder="Your answer…")
        answers[question] = answer
    col1, col2 = st.columns([3, 1])
    with col1:
        if st.button("Continue with this context", use_container_width=True, type="primary"):
            clarifications = "; ".join([f"{q}: {a}" for q, a in answers.items() if a])
            enriched = f"{pending['instruction']}. Clarifications: {clarifications}"
            st.session_state.pending_clarification = None
            with st.spinner(""):
                result = handle_intelligent_edit(enriched)
                if result and result.successful:
                    st.rerun()
    with col2:
        if st.button("Skip", use_container_width=True):
            st.session_state.pending_clarification = None
            with st.spinner(""):
                result = handle_intelligent_edit(pending["instruction"])
                if result and result.successful:
                    st.rerun()


# ============================================================================
# FILE SUGGESTIONS BANNER
# ============================================================================

def render_file_suggestions():
    suggestions = st.session_state.get("pending_file_suggestions", [])
    if not suggestions:
        return
    items = "".join(
        f'<div class="file-sug-item">📎 <strong>{s["file"]}</strong> — '
        f'{", ".join(s["matched_terms"][:3])}</div>'
        for s in suggestions[:3]
    )
    st.markdown(
        f'<div class="file-sug-banner">'
        f'<span class="file-sug-head">💡 File references available</span>{items}</div>',
        unsafe_allow_html=True,
    )
    if st.button("Dismiss", key="dismiss_file_suggestions"):
        st.session_state.pending_file_suggestions = []
        st.rerun()


# ============================================================================
# SMART ROLLBACK
# ============================================================================

def smart_rollback(target_description: str = "") -> bool:
    workspace = st.session_state.workspace
    versions = workspace.version_history
    if not versions:
        return False
    if not target_description:
        after_versions = [v for v in versions if v["description"].startswith("After:")]
        if len(after_versions) >= 2:
            workspace.restore_version(after_versions[-2]["id"])
            return True
        return False
    keyword = target_description.lower()
    for v in reversed(versions):
        if keyword in v["description"].lower():
            workspace.restore_version(v["id"])
            return True
    return False


# ============================================================================
# EDIT PLAN DISPLAY
# ============================================================================

def render_edit_plan(plan) -> None:
    if not plan:
        return
    with st.expander("View edit plan", expanded=False):
        st.markdown(f"**Strategy:** {plan.strategy}")
        for step in plan.steps:
            st.markdown(f"- {step}")
        if plan.rationale:
            st.caption(plan.rationale)


# ============================================================================
# CONVERSATION DASHBOARD
# ============================================================================

def render_conversation_dashboard():
    core = st.session_state.get("intelligence_core")
    if not core:
        return
    conv = core["conversation_manager"]
    if not conv.turns:
        st.caption("No edits yet.")
        return
    evolution = conv.get_document_evolution()
    intent = conv.intent_summary or ""
    if intent:
        st.caption(f"Goal: {intent}")
    for i, turn in enumerate(evolution[-6:], 1):
        changes = turn.get("changes", {})
        net = changes.get("net_change", 0)
        color = "#2e7d32" if net >= 0 else "#c62828"
        label = turn["query"][:38] + ("…" if len(turn["query"]) > 38 else "")
        st.markdown(
            f"<div class='timeline-row'><span class='tl-num'>{i}</span>"
            f"<span class='tl-label'>{label}</span>"
            f"<span class='tl-delta' style='color:{color}'>{net:+d}w</span></div>",
            unsafe_allow_html=True,
        )
    st.markdown("---")
    col_a, col_b = st.columns([3, 1])
    with col_a:
        rollback_kw = st.text_input(
            "Rollback keyword", placeholder="e.g. formal", key="rollback_kw",
            label_visibility="collapsed"
        )
    with col_b:
        if st.button("Undo", use_container_width=True, help="Undo / rollback"):
            kw = rollback_kw.strip() if rollback_kw.strip() else ""
            if smart_rollback(kw):
                st.success("Rolled back")
                st.rerun()
            else:
                st.warning("Nothing found")


# ============================================================================
# SESSION STATE INIT
# ============================================================================

def init_session_state():
    if "workspace" not in st.session_state:
        st.session_state.workspace = DocumentWorkspace()
    if "chat_history" not in st.session_state:
        st.session_state.chat_history = []
    if "uploaded_files" not in st.session_state:
        st.session_state.uploaded_files = {}
    if "show_ppt_download" not in st.session_state:
        st.session_state.show_ppt_download = False
        st.session_state.ppt_data = None
        st.session_state.ppt_topic = ""
    if "show_word_download" not in st.session_state:
        st.session_state.show_word_download = False
        st.session_state.word_data = None
        st.session_state.word_topic = ""
    if "show_excel_download" not in st.session_state:
        st.session_state.show_excel_download = False
        st.session_state.excel_data = None
        st.session_state.excel_topic = ""
    if "intelligence_core" not in st.session_state:
        st.session_state.intelligence_core = None
    if "doc_profile_cache" not in st.session_state:
        st.session_state.doc_profile_cache = None
    if "last_model_used" not in st.session_state:
        st.session_state.last_model_used = None
    if "last_analysis" not in st.session_state:
        st.session_state.last_analysis = None
    if "pending_clarification" not in st.session_state:
        st.session_state.pending_clarification = None
    if "pending_file_suggestions" not in st.session_state:
        st.session_state.pending_file_suggestions = []
    if "selected_ref_files" not in st.session_state:
        st.session_state.selected_ref_files = []
    if "artifact" not in st.session_state:
        st.session_state.artifact = None
    if "sidebar_view" not in st.session_state:
        st.session_state.sidebar_view = "chat"


# ============================================================================
# CSS — Claude-inspired dark design system
# ============================================================================

def apply_custom_css():
    st.markdown("""
    <style>
    @import url('https://fonts.googleapis.com/css2?family=DM+Sans:wght@400;500;600&family=DM+Mono&display=swap');

    :root {
        --bg:           #212121;
        --bg-sidebar:   #171717;
        --bg-input:     #2f2f2f;
        --bg-card:      #2a2a2a;
        --bg-artifact:  #1a1a1a;
        --border:       #383838;
        --border-light: #484848;
        --text:         #ececec;
        --text-muted:   #9a9a9a;
        --text-dim:     #555;
        --accent:       #cc785c;
        --accent-soft:  rgba(204,120,92,0.14);
        --radius:       12px;
        --radius-sm:    8px;
        --font:         'DM Sans', 'ui-sans-serif', system-ui, sans-serif;
        --font-mono:    'DM Mono', 'ui-monospace', monospace;
    }

    html, body,
    [data-testid="stAppViewContainer"],
    [data-testid="stMain"] {
        background: var(--bg) !important;
        color: var(--text) !important;
        font-family: var(--font) !important;
    }

    #MainMenu, footer, header,
    [data-testid="stToolbar"],
    [data-testid="stDecoration"],
    [data-testid="collapsedControl"] { display: none !important; }

    /* Sidebar */
    [data-testid="stSidebar"] {
        background: var(--bg-sidebar) !important;
        border-right: 1px solid var(--border) !important;
    }
    [data-testid="stSidebar"] > div { padding: 0 !important; }

    /* Sidebar logo */
    .sb-logo {
        display: flex; align-items: center; gap: 10px;
        padding: 18px 16px 14px;
        border-bottom: 1px solid var(--border);
        margin-bottom: 4px;
    }
    .sb-logo-icon {
        width: 30px; height: 30px; background: var(--accent);
        border-radius: 7px; display: flex; align-items: center;
        justify-content: center; font-size: 15px; font-weight: 700; color: #fff;
        flex-shrink: 0;
    }
    .sb-logo-name { font-size: 15px; font-weight: 600; color: var(--text); letter-spacing: -.2px; }

    /* Stat rows */
    .stat-row {
        display: flex; justify-content: space-between;
        padding: 5px 0; border-bottom: 1px solid var(--border);
        font-size: 12px;
    }
    .stat-key { color: var(--text-dim); }
    .stat-val { color: var(--text-muted); font-weight: 500; }
    .model-dot { width: 6px; height: 6px; border-radius: 50%;
        display: inline-block; margin-right: 4px; }

    /* Timeline */
    .timeline-row {
        display: flex; align-items: center; gap: 8px;
        padding: 4px 0; font-size: 12px;
    }
    .tl-num  { color: var(--text-dim); font-size: 10.5px; width: 14px; flex-shrink: 0; }
    .tl-label { color: var(--text-muted); flex: 1; overflow: hidden;
        text-overflow: ellipsis; white-space: nowrap; }
    .tl-delta { font-size: 11px; font-weight: 600; flex-shrink: 0; }

    /* Section labels */
    .sec-label {
        font-size: 10.5px; font-weight: 600; color: var(--text-dim);
        text-transform: uppercase; letter-spacing: .6px;
        padding: 12px 0 5px;
    }

    /* Artifact panel */
    .artifact-panel {
        background: var(--bg-artifact); border: 1px solid var(--border);
        border-radius: var(--radius); overflow: hidden; margin-bottom: 16px;
        animation: fadeUp .2s ease;
    }
    @keyframes fadeUp { from { opacity:0; transform:translateY(6px); } to { opacity:1; transform:none; } }

    .artifact-header {
        display: flex; align-items: center; justify-content: space-between;
        padding: 10px 14px; border-bottom: 1px solid var(--border);
        background: var(--bg-card);
    }
    .artifact-badge {
        display: flex; align-items: center; gap: 8px;
        font-size: 12px; font-weight: 600; color: var(--text-muted);
        text-transform: uppercase; letter-spacing: .5px;
    }
    .artifact-icon {
        width: 22px; height: 22px; border-radius: 5px;
        display: flex; align-items: center; justify-content: center; font-size: 12px;
    }
    .art-word  { background: #1565c0; }
    .art-pptx  { background: #b71c1c; }
    .art-excel { background: #1b5e20; }
    .artifact-title { font-size: 13px; font-weight: 500; color: var(--text); }
    .artifact-preview {
        padding: 14px 16px; font-family: var(--font-mono); font-size: 12.5px;
        color: var(--text-muted); max-height: 180px; overflow-y: auto;
        white-space: pre-wrap; line-height: 1.6;
    }
    .artifact-preview::-webkit-scrollbar { width: 3px; }
    .artifact-preview::-webkit-scrollbar-thumb { background: var(--border); }

    /* Edit result card */
    .edit-card {
        background: var(--bg-card); border: 1px solid var(--border);
        border-left: 3px solid var(--accent); border-radius: var(--radius-sm);
        padding: 12px 14px; margin: 10px 0;
    }
    .edit-card-title {
        font-weight: 600; color: var(--text); margin-bottom: 10px;
        font-size: 13px;
    }
    .edit-stats { display: flex; gap: 18px; flex-wrap: wrap; margin-bottom: 8px; }
    .edit-stat { display: flex; flex-direction: column; gap: 1px; }
    .edit-stat-val {
        font-size: 17px; font-weight: 700; color: var(--text);
        font-variant-numeric: tabular-nums;
    }
    .esv-green { color: #4caf50; }
    .esv-red   { color: #ef5350; }
    .edit-stat-label { font-size: 10px; color: var(--text-muted);
        text-transform: uppercase; letter-spacing: .5px; }
    .edit-reason { color: var(--text-muted); font-size: 12.5px;
        line-height: 1.5; margin-bottom: 6px; }
    .edit-meta { color: var(--text-dim); font-size: 11px; }

    /* Clarification */
    .clarif-box {
        background: var(--bg-card); border: 1px solid var(--border-light);
        border-radius: var(--radius-sm); padding: 12px 14px; margin-bottom: 12px;
    }
    .clarif-title { font-weight: 600; font-size: 13.5px; color: var(--text); margin-bottom: 4px; }
    .clarif-sub   { font-size: 12.5px; color: var(--text-muted); }

    /* File suggestion */
    .file-sug-banner {
        background: var(--accent-soft); border: 1px solid rgba(204,120,92,.3);
        border-radius: var(--radius-sm); padding: 10px 14px; margin-bottom: 12px;
    }
    .file-sug-head { font-weight: 600; color: var(--text); display: block; margin-bottom: 5px; }
    .file-sug-item { color: var(--text-muted); margin-bottom: 3px; font-size: 12.5px; }

    /* Footnote */
    .footnote { font-size: 11px; color: var(--text-dim); margin-top: 4px; }

    /* Welcome screen */
    .welcome-icon {
        width: 56px; height: 56px; background: var(--accent);
        border-radius: 14px; display: inline-flex; align-items: center;
        justify-content: center; font-size: 26px; font-weight: 700;
        color: #fff; margin-bottom: 18px;
    }

    /* ── Streamlit overrides ── */
    [data-testid="stTextArea"] textarea,
    [data-testid="stTextInput"] input {
        background: var(--bg-input) !important;
        border: 1px solid var(--border) !important;
        border-radius: var(--radius-sm) !important;
        color: var(--text) !important;
        font-family: var(--font) !important;
        font-size: 14px !important;
    }
    [data-testid="stTextArea"] textarea:focus,
    [data-testid="stTextInput"] input:focus {
        border-color: var(--accent) !important;
        box-shadow: 0 0 0 2px var(--accent-soft) !important;
        outline: none !important;
    }
    [data-testid="stChatInput"] {
        background: var(--bg-input) !important;
        border: 1px solid var(--border) !important;
        border-radius: var(--radius) !important;
    }
    [data-testid="stChatInput"]:focus-within {
        border-color: var(--accent) !important;
        box-shadow: 0 0 0 2px var(--accent-soft) !important;
    }
    [data-testid="stChatInput"] textarea {
        background: transparent !important;
        color: var(--text) !important;
        font-family: var(--font) !important;
        font-size: 14px !important;
    }
    [data-testid="stChatMessage"] {
        background: none !important; border: none !important; padding: 2px 0 !important;
    }
    .stButton > button {
        background: var(--bg-card) !important;
        border: 1px solid var(--border) !important;
        color: var(--text) !important;
        border-radius: var(--radius-sm) !important;
        font-family: var(--font) !important;
        font-size: 13px !important; font-weight: 500 !important;
        transition: background .15s, border-color .15s !important;
        padding: 6px 12px !important;
    }
    .stButton > button:hover {
        background: rgba(255,255,255,.08) !important;
        border-color: var(--border-light) !important;
    }
    .stButton > button[kind="primary"] {
        background: var(--accent) !important;
        border-color: var(--accent) !important;
        color: #fff !important;
    }
    .stButton > button[kind="primary"]:hover { background: #b86b4f !important; }
    .stDownloadButton > button {
        background: var(--accent) !important; border: none !important;
        color: #fff !important; border-radius: var(--radius-sm) !important;
        font-family: var(--font) !important; font-weight: 600 !important;
        font-size: 13px !important; padding: 8px 18px !important;
        width: 100% !important; transition: background .15s !important;
    }
    .stDownloadButton > button:hover { background: #b86b4f !important; }
    [data-testid="stExpander"] {
        background: var(--bg-card) !important;
        border: 1px solid var(--border) !important;
        border-radius: var(--radius-sm) !important;
    }
    [data-testid="stExpander"] summary {
        color: var(--text-muted) !important; font-size: 13px !important;
    }
    [data-testid="stMetric"] {
        background: var(--bg-card) !important;
        border: 1px solid var(--border) !important;
        border-radius: var(--radius-sm) !important;
        padding: 10px !important;
    }
    [data-testid="stMetricLabel"] { color: var(--text-muted) !important; font-size: 11px !important; }
    [data-testid="stMetricValue"] { color: var(--text) !important; font-size: 20px !important; }
    [data-testid="stAlert"] {
        background: var(--bg-card) !important;
        border-radius: var(--radius-sm) !important; font-size: 13px !important;
    }
    [data-testid="stMultiSelect"] div[data-baseweb="select"] {
        background: var(--bg-input) !important; border-color: var(--border) !important;
    }
    [data-testid="stFileUploader"] {
        background: var(--bg-card) !important;
        border: 1px dashed var(--border) !important;
        border-radius: var(--radius-sm) !important;
    }
    * { scrollbar-width: thin; scrollbar-color: var(--border) transparent; }
    ::-webkit-scrollbar { width: 4px; height: 4px; }
    ::-webkit-scrollbar-thumb { background: var(--border); border-radius: 2px; }
    </style>
    """, unsafe_allow_html=True)

# ============================================================================
# SIDEBAR
# ============================================================================

def render_sidebar():
    with st.sidebar:
        st.markdown(
            '<div class="sb-logo">'
            '<div class="sb-logo-icon">M</div>'
            '<div class="sb-logo-name">MozeAI</div>'
            '</div>',
            unsafe_allow_html=True,
        )

        if st.button("＋  New document", key="new_doc_btn", use_container_width=True):
            st.session_state.workspace.current_document["content"] = ""
            st.session_state.workspace.current_document["title"] = "Untitled"
            st.session_state.workspace.save_version("New document")
            st.session_state.chat_history = []
            st.session_state.artifact = None
            st.session_state.intelligence_core["conversation_manager"].clear()
            st.rerun()

        view = st.radio(
            "nav", ["Chat", "Files", "History"],
            key="sidebar_view_radio",
            label_visibility="collapsed",
            horizontal=True,
        )

        if view == "Chat":
            _sidebar_chat_tools()
        elif view == "Files":
            _sidebar_files()
        else:
            _sidebar_history()


def _sidebar_chat_tools():
    core = st.session_state.get("intelligence_core")

    # Intent banner
    if core:
        conv = core["conversation_manager"]
        if conv.turns:
            intent = conv.intent_summary or ""
            if intent:
                st.markdown(
                    f'<div style="font-size:12px;color:var(--text-muted);background:var(--bg-card);'
                    f'border:1px solid var(--border);border-radius:8px;padding:8px 10px;margin:6px 0 10px;">'
                    f'🎯 {intent}</div>',
                    unsafe_allow_html=True,
                )

    st.markdown('<div class="sec-label">Quick edits</div>', unsafe_allow_html=True)
    quick = [
        ("✨ Improve writing",  "improve this document"),
        ("🎓 Make academic",    "make this academic"),
        ("📝 Summarise",        "summarize this document"),
        ("🔧 Fix grammar",      "fix grammar"),
        ("📢 Make persuasive",  "make this persuasive"),
        ("🪄 Simplify",         "simplify this document"),
    ]
    for label, instr in quick:
        if st.button(label, use_container_width=True, key=f"qa_{label[:8]}"):
            with st.spinner(""):
                result = handle_intelligent_edit(instr)
                if result and result.successful:
                    st.rerun()

    st.markdown('<div class="sec-label" style="margin-top:14px">Custom instruction</div>', unsafe_allow_html=True)
    custom = st.text_area(
        "", placeholder="e.g. Rewrite for a 10-year-old…",
        height=72, key="custom_instr", label_visibility="collapsed",
    )
    if st.button("Apply", use_container_width=True, type="primary", key="apply_custom"):
        if custom.strip():
            with st.spinner(""):
                result = handle_intelligent_edit(custom.strip())
                if result and result.successful:
                    st.rerun()

    if core:
        meta = st.session_state.workspace.current_document["metadata"]
        lm = st.session_state.get("last_model_used", MODEL_PRIORITY[0])
        is_primary = lm == MODEL_PRIORITY[0]
        dot_bg = "#4caf50" if is_primary else "#ff9800"
        lm_short = lm.replace("-versatile", "")
        st.markdown(
            f'<div style="margin-top:16px">'
            f'<div class="stat-row"><span class="stat-key">Words</span>'
            f'<span class="stat-val">{meta["word_count"]}</span></div>'
            f'<div class="stat-row"><span class="stat-key">Versions</span>'
            f'<span class="stat-val">{len(st.session_state.workspace.version_history)}</span></div>'
            f'<div class="stat-row"><span class="stat-key">Model</span>'
            f'<span class="stat-val"><span class="model-dot" style="background:{dot_bg}"></span>{lm_short}</span></div>'
            f'</div>',
            unsafe_allow_html=True,
        )


def _sidebar_files():
    st.markdown('<div class="sec-label">Upload files</div>', unsafe_allow_html=True)
    uploaded = st.file_uploader(
        "", type=["pdf", "docx", "txt", "csv", "json"],
        accept_multiple_files=True, key="file_uploader",
        label_visibility="collapsed",
    )
    if uploaded:
        for file in uploaded:
            if file.name not in st.session_state.uploaded_files:
                content = process_uploaded_file(file)
                if content and not content.startswith("Error"):
                    st.session_state.uploaded_files[file.name] = content
                    st.session_state.intelligence_core["file_accumulator"].add_file(
                        file.name, file.type, content, {}
                    )
                    st.toast(f"Loaded {file.name}", icon="✅")

    all_files = list(st.session_state.uploaded_files.keys())
    if all_files:
        st.markdown('<div class="sec-label" style="margin-top:10px">Reference in next edit</div>',
                    unsafe_allow_html=True)
        st.multiselect(
            "", options=all_files, default=[],
            key="selected_ref_files", label_visibility="collapsed",
        )
        sel = st.session_state.get("selected_ref_files", [])
        if sel:
            st.caption(f"{len(sel)} file(s) selected")
    else:
        st.markdown(
            '<div style="color:var(--text-dim);font-size:12.5px;text-align:center;padding:24px 0;">'
            'No files uploaded yet</div>',
            unsafe_allow_html=True,
        )


def _sidebar_history():
    st.markdown('<div class="sec-label">Edit timeline</div>', unsafe_allow_html=True)
    render_conversation_dashboard()

    versions = st.session_state.workspace.version_history
    if versions:
        st.markdown('<div class="sec-label" style="margin-top:12px">Version history</div>',
                    unsafe_allow_html=True)
        for v in reversed(versions[-8:]):
            desc = v["description"][:38]
            wc = v["metadata"].get("word_count", 0)
            st.markdown(
                f'<div style="padding:5px 0;border-bottom:1px solid var(--border);'
                f'font-size:12px;color:var(--text-muted);display:flex;'
                f'justify-content:space-between;">'
                f'<span>{desc}</span><span style="color:var(--text-dim)">{wc}w</span></div>',
                unsafe_allow_html=True,
            )


# ============================================================================
# ARTIFACT PANEL  — shows generated docs like Claude's artifact pane
# ============================================================================

def render_artifact_panel():
    artifact = st.session_state.get("artifact")
    if not artifact:
        return

    type_conf = {
        "word":  ("📄", "art-word",  "Word Document"),
        "pptx":  ("📊", "art-pptx",  "PowerPoint"),
        "excel": ("📈", "art-excel", "Spreadsheet"),
    }
    icon, css_cls, type_label = type_conf.get(
        artifact["type"], ("📄", "art-word", "Document")
    )

    preview_safe = (
        artifact.get("preview", "")
        .replace("&", "&amp;").replace("<", "&lt;").replace(">", "&gt;")[:600]
    )

    st.markdown(
        f'<div class="artifact-panel">'
        f'  <div class="artifact-header">'
        f'    <div class="artifact-badge">'
        f'      <div class="artifact-icon {css_cls}">{icon}</div>'
        f'      <span>{type_label}</span>'
        f'    </div>'
        f'    <div class="artifact-title">{artifact["title"]}</div>'
        f'  </div>'
        f'  <div class="artifact-preview">{preview_safe}</div>'
        f'</div>',
        unsafe_allow_html=True,
    )

    col_dl, col_close = st.columns([4, 1])
    with col_dl:
        st.download_button(
            f"⬇  Download {artifact['title']}",
            data=artifact["bytes"],
            file_name=artifact["filename"],
            mime=artifact["mime"],
            use_container_width=True,
        )
    with col_close:
        if st.button("✕", use_container_width=True, key="close_artifact",
                     help="Close"):
            st.session_state.artifact = None
            st.rerun()


def _maybe_show_artifact(response):
    """Convert pending download state into a rich artifact card."""
    if st.session_state.get("show_word_download") and st.session_state.get("word_data"):
        data  = st.session_state.word_data
        topic = st.session_state.word_topic
        try:
            from docx import Document as _D
            buf = data if hasattr(data, "read") else BytesIO(data)
            doc = _D(buf)
            preview = "\n".join(p.text for p in doc.paragraphs if p.text.strip())[:600]
        except Exception:
            preview = f"Word document: {topic}"
        st.session_state.artifact = {
            "type": "word", "title": f"{topic}.docx",
            "filename": f"{topic}.docx",
            "mime": "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
            "bytes": st.session_state.word_data, "preview": preview,
        }
        st.session_state.show_word_download = False

    elif st.session_state.get("show_ppt_download") and st.session_state.get("ppt_data"):
        topic = st.session_state.ppt_topic
        st.session_state.artifact = {
            "type": "pptx", "title": f"{topic}.pptx",
            "filename": f"{topic}.pptx",
            "mime": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
            "bytes": st.session_state.ppt_data,
            "preview": f"PowerPoint: {topic}\nGenerated by MozeAI.",
        }
        st.session_state.show_ppt_download = False

    elif st.session_state.get("show_excel_download") and st.session_state.get("excel_data"):
        topic = st.session_state.excel_topic
        st.session_state.artifact = {
            "type": "excel", "title": f"{topic}.xlsx",
            "filename": f"{topic}.xlsx",
            "mime": "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
            "bytes": st.session_state.excel_data,
            "preview": f"Excel spreadsheet: {topic}\nCreated by MozeAI.",
        }
        st.session_state.show_excel_download = False


# ============================================================================
# DOCUMENT EDITOR
# ============================================================================

def render_document_editor():
    workspace = st.session_state.workspace
    meta = workspace.current_document["metadata"]

    col_title, col_btns = st.columns([3, 2])
    with col_title:
        new_title = st.text_input(
            "", value=workspace.current_document["title"],
            key="doc_title", placeholder="Document title…",
            label_visibility="collapsed",
        )
        if new_title != workspace.current_document["title"]:
            workspace.current_document["title"] = new_title

    with col_btns:
        c1, c2, c3, c4 = st.columns(4)
        with c1:
            if st.button("Analyse", key="btn_analyse", use_container_width=True):
                profile = profile_current_document()
                st.session_state.doc_profile_cache = profile
                tone  = profile.content.get("tone", "—")
                level = profile.content.get("reading_level", "—")
                st.toast(f"Tone: {tone}  ·  {level}", icon="🔍")
        with c2:
            track = workspace.track_changes
            if st.button("Track ✓" if track else "Track", key="btn_track",
                         use_container_width=True):
                workspace.track_changes = not track
                st.rerun()
        with c3:
            if st.button("Stats", key="btn_stats", use_container_width=True):
                st.toast(f"{meta['word_count']} words · {meta['reading_time']} min read",
                         icon="📊")
        with c4:
            if st.button("Clear", key="btn_clear", use_container_width=True):
                workspace.current_document["content"] = ""
                workspace._update_metadata()
                st.rerun()

    content = st.text_area(
        "", value=workspace.current_document["content"],
        height=340, key="doc_editor",
        placeholder="Start writing, or paste your document here…",
        label_visibility="collapsed",
    )
    if content != workspace.current_document["content"]:
        workspace.update_document(content, "Manual edit")

    st.markdown(
        f'<div class="footnote">{meta["word_count"]} words '
        f'· {meta["char_count"]} chars '
        f'· {meta["reading_time"]} min read</div>',
        unsafe_allow_html=True,
    )


# ============================================================================
# CHAT INTERFACE
# ============================================================================

def render_chat_interface():
    render_clarification_dialog()
    render_file_suggestions()

    for role, msg in st.session_state.chat_history[-20:]:
        avatar = "M" if role == "assistant" else "👤"
        with st.chat_message(role, avatar=avatar):
            st.markdown(msg)

    render_artifact_panel()

    query = st.chat_input("Ask anything — edit, generate a Word doc, make a deck…")

    if query and query.strip().lower() in ("undo", "undo last"):
        with st.chat_message("user", avatar="👤"):
            st.markdown(query)
        st.session_state.chat_history.append(("user", query))
        with st.chat_message("assistant", avatar="M"):
            msg = ("↩️ Done — restored the previous version."
                   if smart_rollback() else "⚠️ Nothing to undo.")
            st.markdown(msg)
            st.session_state.chat_history.append(("assistant", msg))
        st.rerun()
        return

    if not query:
        return

    st.session_state.chat_history.append(("user", query))
    with st.chat_message("user", avatar="👤"):
        st.markdown(query)

    q_lower = query.lower().strip()
    edit_keywords = [
        "improve", "rewrite", "summarize", "summarise", "expand", "shorten",
        "fix grammar", "make formal", "make academic", "translate",
        "make persuasive", "simplify",
    ]
    is_edit = (
        any(kw in q_lower for kw in edit_keywords)
        and st.session_state.workspace.current_document["content"]
        and not q_lower.startswith("/")
    )

    with st.chat_message("assistant", avatar="M"):
        if is_edit:
            resp_ph  = st.empty()
            meta_ph  = st.empty()
            buf      = {"text": "", "tok": 0, "t0": time.time()}

            def on_tok(token):
                buf["text"] += token
                buf["tok"]  += 1
                elapsed = max(time.time() - buf["t0"], 0.001)
                tps = round(buf["tok"] / elapsed, 1)
                resp_ph.markdown(buf["text"] + " ▌")
                meta_ph.markdown(
                    f'<div class="footnote">{len(buf["text"].split())} words · {tps} tok/s</div>',
                    unsafe_allow_html=True,
                )

            sel = st.session_state.get("selected_ref_files", []) or None
            response = run_agent(query, stream_callback=on_tok, selected_files=sel)
            meta_ph.empty()

            if response == "__CLARIFICATION_NEEDED__":
                resp_ph.empty()
                st.rerun()
                return

            if isinstance(response, EditResult) and response.successful:
                resp_ph.empty()
                ch = response.changes_made
                old_wc = ch.get("old_word_count", 0)
                new_wc = ch.get("new_word_count", 0)
                net    = ch.get("net_change", 0)
                adds   = ch.get("additions", 0)
                dels   = ch.get("deletions", 0)
                ms     = response.execution_time_ms
                sh     = st.session_state.intelligence_core["streaming_handler"]
                tps    = sh.last_metrics.get("tokens_per_sec", 0)
                ttft   = sh.last_metrics.get("ttft_ms", 0)
                toks   = sh.last_metrics.get("token_count", 0)
                mdl    = sh.last_metrics.get("model", MODEL_PRIORITY[0]).replace("-versatile","")

                add_cls = "esv-green" if adds > 0 else ""
                del_cls = "esv-red"   if dels > 0 else ""
                net_cls = "esv-green" if net >= 0 else "esv-red"
                rsn_html = (f'<div class="edit-reason">💡 {response.reasoning}</div>'
                            if response.reasoning else "")

                st.markdown(f"""
                <div class="edit-card">
                  <div class="edit-card-title">✅ Edit complete</div>
                  <div class="edit-stats">
                    <div class="edit-stat">
                      <div class="edit-stat-val {add_cls}">+{adds}</div>
                      <div class="edit-stat-label">Added</div>
                    </div>
                    <div class="edit-stat">
                      <div class="edit-stat-val {del_cls}">−{dels}</div>
                      <div class="edit-stat-label">Removed</div>
                    </div>
                    <div class="edit-stat">
                      <div class="edit-stat-val {net_cls}">{net:+d}</div>
                      <div class="edit-stat-label">Net</div>
                    </div>
                    <div class="edit-stat">
                      <div class="edit-stat-val">{ms}ms</div>
                      <div class="edit-stat-label">Time</div>
                    </div>
                    <div class="edit-stat">
                      <div class="edit-stat-val">{tps}</div>
                      <div class="edit-stat-label">tok/s</div>
                    </div>
                  </div>
                  {rsn_html}
                  <div class="edit-meta">
                    {old_wc} → {new_wc} words · {toks} tokens · TTFT {ttft}ms · {mdl}
                  </div>
                </div>
                """, unsafe_allow_html=True)

                core = st.session_state.intelligence_core
                try:
                    snap   = {"content": st.session_state.workspace.current_document["content"],
                              "title":   st.session_state.workspace.current_document["title"],
                              "word_count": new_wc}
                    parsed = core["instruction_parser"].parse(query, snap)
                    plan   = core["edit_planner"].plan(parsed, snap, core["conversation_manager"])
                    render_edit_plan(plan)
                except Exception:
                    pass

                st.session_state.chat_history.append((
                    "assistant",
                    f"✅ Edit complete in {ms}ms · {net:+d} words ({old_wc}→{new_wc}) · {tps} tok/s",
                ))

            elif isinstance(response, EditResult) and not response.successful:
                msg = f"⚠️ {response.reasoning}"
                resp_ph.warning(msg)
                st.session_state.chat_history.append(("assistant", msg))

            elif isinstance(response, str):
                resp_ph.markdown(response)
                st.session_state.chat_history.append(("assistant", response))

        else:
            # General / doc-generation queries with streaming
            resp_ph  = st.empty()
            meta_ph  = st.empty()
            buf      = {"text": "", "tok": 0, "t0": time.time()}

            def on_tok_gen(token):
                buf["text"] += token
                buf["tok"]  += 1
                elapsed = max(time.time() - buf["t0"], 0.001)
                tps = round(buf["tok"] / elapsed, 1)
                resp_ph.markdown(buf["text"] + " ▌")
                meta_ph.markdown(f'<div class="footnote">⚡ {tps} tok/s</div>',
                                  unsafe_allow_html=True)

            response = run_agent(query, stream_callback=on_tok_gen)
            meta_ph.empty()
            resp_ph.empty()

            _maybe_show_artifact(response)

            if isinstance(response, str):
                st.markdown(response)
                st.session_state.chat_history.append(("assistant", response))

    st.rerun()


# ============================================================================
# MAIN
# ============================================================================

def main():
    st.set_page_config(
        page_title="MozeAI",
        page_icon="M",
        layout="wide",
        initial_sidebar_state="expanded",
    )

    init_session_state()
    apply_custom_css()

    # Groq client
    groq_api_key = None
    try:
        if "GROQ_API_KEY" in st.secrets:
            groq_api_key = st.secrets["GROQ_API_KEY"]
    except Exception:
        pass
    if not groq_api_key:
        groq_api_key = os.environ.get("GROQ_API_KEY")
    if not groq_api_key:
        st.error("GROQ_API_KEY not found. Set it in Streamlit Secrets or as an env var.")
        st.stop()

    global client
    client = Groq(api_key=groq_api_key)

    # Intelligence core
    if st.session_state.intelligence_core is None:
        sh = StreamingResponseHandler(client)
        st.session_state.intelligence_core = {
            "conversation_manager": ConversationManager(),
            "instruction_parser":   InstructionParser(client),
            "file_accumulator":     FileContextAccumulator(client),
            "edit_planner":         EditPlanner(client),
            "streaming_handler":    sh,
            "document_profiler":    DocumentProfiler(client),
            "contextual_editor":    ContextualEditor(client, sh),
        }

    render_sidebar()

    has_content = bool(st.session_state.workspace.current_document["content"].strip())
    has_history = bool(st.session_state.chat_history)

    if has_content or has_history:
        editor_col, chat_col = st.columns([5, 6], gap="large")
        with editor_col:
            st.markdown('<div style="padding-top:14px">', unsafe_allow_html=True)
            render_document_editor()
            st.markdown("</div>", unsafe_allow_html=True)
        with chat_col:
            st.markdown('<div style="padding-top:14px">', unsafe_allow_html=True)
            render_chat_interface()
            st.markdown("</div>", unsafe_allow_html=True)
    else:
        # Welcome / empty state
        _, centre, _ = st.columns([1, 4, 1])
        with centre:
            st.markdown("""
            <div style="text-align:center;padding:64px 0 32px">
              <div class="welcome-icon">M</div>
              <div style="font-size:24px;font-weight:600;color:var(--text);
                          letter-spacing:-.4px;margin-bottom:8px;">
                How can I help you today?
              </div>
              <div style="font-size:14px;color:var(--text-muted);
                          max-width:400px;margin:0 auto;line-height:1.6;">
                Write or paste a document, upload files, or ask me to create
                a Word doc, PowerPoint, or spreadsheet — just like you would with Claude.
              </div>
            </div>
            """, unsafe_allow_html=True)

            suggestions = [
                "Write a business proposal about renewable energy",
                "Make a PowerPoint on AI trends in 2025",
                "Create an Excel sales tracker template",
                "Summarise this document for executives",
            ]
            c1, c2 = st.columns(2)
            for i, sug in enumerate(suggestions):
                col = c1 if i % 2 == 0 else c2
                with col:
                    if st.button(sug, use_container_width=True, key=f"sug_{i}"):
                        st.session_state.chat_history.append(("user", sug))
                        with st.spinner(""):
                            response = run_agent(sug)
                        _maybe_show_artifact(response)
                        if isinstance(response, str):
                            st.session_state.chat_history.append(("assistant", response))
                        st.rerun()

            render_chat_interface()


if __name__ == "__main__":
    main()
